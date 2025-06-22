# modules/file_processor.py

import streamlit as st
import os
import time
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
import pandas as pd

# Import untuk format dokumen
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Import pembolehubah konfigurasi
from config import UPLOAD_DIR, LOGO_PATH, WATERMARK_TEXT, FONT_DIR

# --- FUNGSI EKSTRAKSI TEKS ---

def extract_text_from_file(uploaded_file_obj):
    """
    Mengekstrak teks dari objek fail yang dimuat naik (Imej, PDF, DOCX, TXT).
    
    Args:
        uploaded_file_obj: Objek fail dari st.file_uploader.

    Returns:
        String teks yang diekstrak, atau None jika gagal atau jenis fail tidak disokong.
    """
    filename = uploaded_file_obj.name
    file_bytes = uploaded_file_obj.getvalue()
    extracted_text = ""
    temp_file_path = None

    try:
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            image = Image.open(uploaded_file_obj)
            extracted_text = pytesseract.image_to_string(image)
            if not extracted_text.strip():
                st.info(f"Tiada teks dapat diekstrak dari imej '{filename}' menggunakan OCR.")
        
        elif filename.lower().endswith(".txt"):
            extracted_text = file_bytes.decode('utf-8', errors='ignore')

        elif filename.lower().endswith(".docx"):
            # Simpan fail sementara untuk dibaca oleh python-docx
            temp_file_path = os.path.join(UPLOAD_DIR, f"temp_{int(time.time())}_{filename}")
            with open(temp_file_path, "wb") as f:
                f.write(file_bytes)
            doc = Document(temp_file_path)
            extracted_text = "\n".join([para.text for para in doc.paragraphs])

        elif filename.lower().endswith(".pdf"):
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                extracted_text = "".join(page.get_text() for page in doc)
        
        else:
            st.warning(f"Jenis fail '{filename}' tidak disokong untuk ekstraksi teks.")
            return None
            
        return extracted_text.strip() if isinstance(extracted_text, str) else None

    except Exception as e:
        st.error(f"Ralat semasa memproses fail '{filename}': {e}")
        return None
    
    finally:
        # Pastikan fail sementara dipadam
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.remove(temp_file_path)
            except OSError as e_remove:
                st.warning(f"Gagal memadam fail sementara '{temp_file_path}': {e_remove}")


# --- FUNGSI PEMFORMATAN & EKSPORT ---

def format_conversation_text(chat_history, include_user=True, include_assistant=True):
    """Memformat sejarah perbualan menjadi satu string teks yang boleh dibaca."""
    lines = []
    for msg in chat_history:
        role = msg.get("role", "unknown").capitalize()
        content = msg.get("content", "").strip()
        
        if (msg["role"] == "user" and include_user) or (msg["role"] == "assistant" and include_assistant):
            lines.append(f"{role}: {content}")
            if msg["role"] == "assistant" and msg.get("thinking_process"):
                thinking_process = msg["thinking_process"].strip()
                lines.append(f"  [Proses Pemikiran AI]:\n  {thinking_process}\n")
    
    return "\n\n".join(lines)

def save_to_word(text_content, full_path):
    """Menyimpan kandungan teks ke fail .docx."""
    try:
        doc = Document()
        # Tambah logo jika wujud
        if LOGO_PATH and os.path.exists(LOGO_PATH):
            try:
                p = doc.add_paragraph()
                p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                p.add_run().add_picture(LOGO_PATH, width=DocxInches(1.5))
                doc.add_paragraph() # Jarak
            except Exception as e:
                st.warning(f"Gagal menambah logo pada Word: {e}")

        # Tambah kandungan utama
        for para_block in text_content.split("\n\n"):
            doc.add_paragraph(para_block.strip())
        
        doc.save(full_path)
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke Word: {e}")
        return False

def save_to_pdf(text_content, full_path):
    """Menyimpan kandungan teks ke fail .pdf."""
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Tetapan Fon Unicode (Penting untuk menyokong pelbagai aksara)
        font_path = os.path.join(FONT_DIR, "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdf.add_font("DejaVu", "", font_path, uni=True)
            pdf.set_font("DejaVu", size=12)
        else:
            pdf.set_font("Arial", size=12)
            if "DejaVu" in str(font_path): # Hanya tunjuk amaran jika fon sepatutnya ada
                 st.warning("Fail fon 'DejaVuSans.ttf' tidak ditemui. Menggunakan fon lalai.")

        # Tambah logo jika wujud
        if LOGO_PATH and os.path.exists(LOGO_PATH):
            try:
                page_width = pdf.w - 2 * pdf.l_margin
                x_logo = (page_width - 30) / 2 + pdf.l_margin
                pdf.image(LOGO_PATH, x=x_logo, y=10, w=30)
                pdf.ln(25) # Jarak selepas logo
            except Exception as e:
                st.warning(f"Gagal menambah logo pada PDF: {e}")

        # Tambah kandungan utama
        pdf.multi_cell(0, 10, text_content)
        
        pdf.output(full_path)
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke PDF: {e}")
        return False

def save_to_txt(text_content, full_path):
    """Menyimpan kandungan teks ke fail .txt."""
    try:
        with open(full_path, "w", encoding="utf-8") as f:
            f.write(text_content)
        return True
    except IOError as e:
        st.error(f"Gagal menyimpan ke Teks: {e}")
        return False

def save_to_excel(chat_history, full_path):
    """Menyimpan sejarah perbualan ke fail .xlsx."""
    try:
        data = [
            {
                "Peranan": msg.get("role", "").capitalize(),
                "Mesej": msg.get("content", ""),
                "Proses Pemikiran": msg.get("thinking_process", "")
            }
            for msg in chat_history
        ]
        if not data:
            st.warning("Tiada data untuk dieksport ke Excel.")
            return False
            
        df = pd.DataFrame(data)
        df.to_excel(full_path, index=False, engine='openpyxl')
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke Excel: {e}")
        return False

def save_to_pptx(chat_history, full_path):
    """Menyimpan sejarah perbualan ke fail .pptx, satu mesej per slaid."""
    try:
        prs = Presentation()
        # Gunakan layout 'Title and Content' atau 'Blank'
        slide_layout = prs.slide_layouts[5] 

        for msg in chat_history:
            slide = prs.slides.add_slide(slide_layout)
            
            # Tajuk slaid (Peranan)
            title = slide.shapes.title
            title.text = f"Peranan: {msg.get('role', '').capitalize()}"

            # Kandungan slaid
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Kosongkan perenggan sedia ada
            
            p_content = tf.add_paragraph()
            p_content.text = msg.get("content", "(Tiada kandungan)")
            p_content.font.size = PptxPt(18)
            
            # Tambah proses pemikiran jika ada
            if msg.get("thinking_process"):
                tf.add_paragraph() # Jarak
                p_think_header = tf.add_paragraph()
                p_think_header.text = "Proses Pemikiran AI:"
                p_think_header.font.bold = True
                p_think_header.font.size = PptxPt(14)
                
                p_think_content = tf.add_paragraph()
                p_think_content.text = msg.get("thinking_process")
                p_think_content.font.size = PptxPt(12)
                p_think_content.level = 1

        prs.save(full_path)
        return True
    except Exception as e:
        st.error(f"Gagal menyimpan ke PowerPoint: {e}")
        return False
